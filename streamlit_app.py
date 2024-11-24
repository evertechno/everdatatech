import streamlit as st
import os
from tempfile import NamedTemporaryFile
from fpdf import FPDF
import markdown
from PIL import Image
import json
import xml.etree.ElementTree as ET
import pyqrcode
import barcode
from barcode.writer import ImageWriter
import pytz
from datetime import datetime
from io import BytesIO
import pandas as pd
import pyttsx3
import speech_recognition as sr
from moviepy.editor import VideoFileClip
import ebooklib
from ebooklib import epub
import shutil
import zipfile
from pptx import Presentation
import base64
import csv
import io
import docx
import re
import shutil
from collections import Counter
from urllib.parse import urlparse
import hashlib
import random
import string
import difflib
import itertools
import socket
from datetime import timedelta

# Streamlit App UI Setup
st.title("Advanced File Conversion Platform with AI Assistant")
st.write("Convert various file formats and get assistance from the AI assistant.")

# File upload section
uploaded_file = st.file_uploader("Upload your file", type=["dwg", "rvt", "ai", "fdr", "pptx", "txt", "zip", "jpg", "jpeg", "png", "md", "html", "epub", "json", "xml", "mp3", "mp4", "docx"])

# Conversion options
conversion_type = st.selectbox(
    "Select the conversion format",
    [
        "DWG to PDF", "RVT to DWG", "AI to PDF", "FDR to PDF", "PPT to PDF", "TXT to PDF", 
        "MD to PDF", "HTML to PDF", "EPUB to PDF", "JSON to PDF", "XML to PDF", "Extract ZIP", 
        "Compress Folder", "Image to PDF", "Audio to Text", "Text to Speech", "Video to Audio", 
        "QR Code Generator", "Barcode Generator", "HTML to Text", "CSV to Excel", "Excel to CSV",
        "Markdown to Text", "Text to Markdown", "Text File Cleaner", "Word to PDF", "CSV to JSON",
        "URL to HTML", "HTML to CSV", "Count Word Frequency", "String Similarity", "Random Password Generator",
        "Hash String", "File Merger", "Convert Timezone", "Convert Case", "Create Custom JSON", "Date Difference",
        "Unique Items Finder", "Compress Text", "Markdown to HTML", "Text to JSON", "CSV to TXT"
    ]
)

# File conversion functions
def convert_ppt_to_pdf(input_file, output_file):
    try:
        ppt = Presentation(input_file)
        pdf_writer = FPDF()
        pdf_writer.add_page()
        for slide in ppt.slides:
            pdf_writer.cell(200, 10, txt=slide.shapes.title.text, ln=True)
        pdf_writer.output(output_file)
        return output_file
    except Exception as e:
        return f"Error in PPT to PDF conversion: {str(e)}"

def convert_txt_to_pdf(input_file, output_file):
    try:
        with open(input_file, 'r') as file:
            text = file.read()
        pdf = FPDF()
        pdf.add_page()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.set_font("Arial", size=12)
        pdf.multi_cell(0, 10, text)
        pdf.output(output_file)
        return output_file
    except Exception as e:
        return f"Error in TXT to PDF conversion: {str(e)}"

def convert_md_to_pdf(input_file, output_file):
    try:
        with open(input_file, 'r') as file:
            text = markdown.markdown(file.read())
        pdf = FPDF()
        pdf.add_page()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.set_font("Arial", size=12)
        pdf.multi_cell(0, 10, text)
        pdf.output(output_file)
        return output_file
    except Exception as e:
        return f"Error in Markdown to PDF conversion: {str(e)}"

def convert_html_to_pdf(input_file, output_file):
    try:
        html_content = open(input_file, 'r').read()
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        pdf.multi_cell(0, 10, html_content)
        pdf.output(output_file)
        return output_file
    except Exception as e:
        return f"Error in HTML to PDF conversion: {str(e)}"

def convert_epub_to_pdf(input_file, output_file):
    try:
        book = epub.read_epub(input_file)
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                pdf.multi_cell(0, 10, item.content.decode())
        pdf.output(output_file)
        return output_file
    except Exception as e:
        return f"Error in EPUB to PDF conversion: {str(e)}"

def convert_json_to_pdf(input_file, output_file):
    try:
        with open(input_file, 'r') as f:
            data = json.load(f)
        pdf = FPDF()
        pdf.add_page()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.set_font("Arial", size=12)
        pdf.multi_cell(0, 10, json.dumps(data, indent=4))
        pdf.output(output_file)
        return output_file
    except Exception as e:
        return f"Error in JSON to PDF conversion: {str(e)}"

def convert_xml_to_pdf(input_file, output_file):
    try:
        tree = ET.parse(input_file)
        root = tree.getroot()
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        pdf.multi_cell(0, 10, ET.tostring(root, encoding='unicode'))
        pdf.output(output_file)
        return output_file
    except Exception as e:
        return f"Error in XML to PDF conversion: {str(e)}"

def generate_qr_code(data, output_file):
    try:
        qr = pyqrcode.create(data)
        qr.png(output_file, scale=6)
        return output_file
    except Exception as e:
        return f"Error generating QR code: {str(e)}"

def generate_barcode(data, output_file):
    try:
        barcode_obj = barcode.get('ean13', data, writer=ImageWriter())
        barcode_obj.save(output_file)
        return output_file
    except Exception as e:
        return f"Error generating barcode: {str(e)}"

def convert_image_to_pdf(input_file, output_file):
    try:
        image = Image.open(input_file)
        image.convert('RGB').save(output_file, "PDF")
        return output_file
    except Exception as e:
        return f"Error in Image to PDF conversion: {str(e)}"

def convert_audio_to_text(input_file):
    try:
        recognizer = sr.Recognizer()
        audio = sr.AudioFile(input_file)
        with audio as source:
            audio_data = recognizer.record(source)
        return recognizer.recognize_google(audio_data)
    except Exception as e:
        return f"Error in audio to text conversion: {str(e)}"

def text_to_speech(input_text, output_audio):
    try:
        engine = pyttsx3.init()
        engine.save_to_file(input_text, output_audio)
        engine.runAndWait()
        return output_audio
    except Exception as e:
        return f"Error in text to speech conversion: {str(e)}"

def convert_video_to_audio(input_file, output_audio):
    try:
        video_clip = VideoFileClip(input_file)
        video_clip.audio.write_audiofile(output_audio)
        return output_audio
    except Exception as e:
        return f"Error in video to audio conversion: {str(e)}"

def extract_zip(input_file, output_folder):
    try:
        with zipfile.ZipFile(input_file, 'r') as zip_ref:
            zip_ref.extractall(output_folder)
        return f"Extracted to {output_folder}"
    except Exception as e:
        return f"Error in zip extraction: {str(e)}"

def compress_folder(input_folder, output_file):
    try:
        shutil.make_archive(output_file, 'zip', input_folder)
        return f"Folder compressed to {output_file}.zip"
    except Exception as e:
        return f"Error in folder compression: {str(e)}"

# Additional Functions
def convert_docx_to_pdf(input_file, output_file):
    try:
        from docx2txt import process
        process(input_file, output_file)
        return output_file
    except Exception as e:
        return f"Error in DOCX to PDF conversion: {str(e)}"

def csv_to_json(input_file, output_file):
    try:
        df = pd.read_csv(input_file)
        df.to_json(output_file, orient="records", lines=True)
        return output_file
    except Exception as e:
        return f"Error in CSV to JSON conversion: {str(e)}"

def url_to_html(input_url, output_file):
    try:
        from requests import get
        response = get(input_url)
        with open(output_file, "w") as file:
            file.write(response.text)
        return output_file
    except Exception as e:
        return f"Error in URL to HTML conversion: {str(e)}"

def count_word_frequency(input_text):
    try:
        words = input_text.split()
        return dict(Counter(words))
    except Exception as e:
        return f"Error in counting word frequency: {str(e)}"

def string_similarity(input_string1, input_string2):
    try:
        return difflib.SequenceMatcher(None, input_string1, input_string2).ratio()
    except Exception as e:
        return f"Error in string similarity: {str(e)}"

def generate_random_password(length=12):
    try:
        chars = string.ascii_letters + string.digits + string.punctuation
        return ''.join(random.choice(chars) for _ in range(length))
    except Exception as e:
        return f"Error in random password generation: {str(e)}"

def hash_string(input_string):
    try:
        return hashlib.sha256(input_string.encode()).hexdigest()
    except Exception as e:
        return f"Error in hashing string: {str(e)}"

def merge_files(input_files, output_file):
    try:
        with open(output_file, 'wb') as f_out:
            for input_file in input_files:
                with open(input_file, 'rb') as f_in:
                    shutil.copyfileobj(f_in, f_out)
        return output_file
    except Exception as e:
        return f"Error in file merging: {str(e)}"

def convert_timezone(input_datetime, from_tz, to_tz):
    try:
        local_time = pytz.timezone(from_tz).localize(input_datetime)
        new_time = local_time.astimezone(pytz.timezone(to_tz))
        return new_time
    except Exception as e:
        return f"Error in timezone conversion: {str(e)}"

def convert_case(input_string, case_type="upper"):
    try:
        if case_type == "upper":
            return input_string.upper()
        elif case_type == "lower":
            return input_string.lower()
        elif case_type == "title":
            return input_string.title()
        else:
            return input_string
    except Exception as e:
        return f"Error in case conversion: {str(e)}"

def create_custom_json(data):
    try:
        return json.dumps(data, indent=4)
    except Exception as e:
        return f"Error in creating custom JSON: {str(e)}"

def date_difference(date1, date2):
    try:
        delta = date2 - date1
        return str(delta)
    except Exception as e:
        return f"Error in date difference calculation: {str(e)}"

def find_unique_items(input_list):
    try:
        return list(set(input_list))
    except Exception as e:
        return f"Error in finding unique items: {str(e)}"

def compress_text(input_text):
    try:
        return zlib.compress(input_text.encode())
    except Exception as e:
        return f"Error in compressing text: {str(e)}"

def markdown_to_html(input_file):
    try:
        with open(input_file, 'r') as file:
            text = file.read()
        return markdown.markdown(text)
    except Exception as e:
        return f"Error in converting markdown to HTML: {str(e)}"

def text_to_json(input_text):
    try:
        data = {"text": input_text}
        return json.dumps(data, indent=4)
    except Exception as e:
        return f"Error in converting text to JSON: {str(e)}"

def csv_to_txt(input_file, output_file):
    try:
        df = pd.read_csv(input_file)
        df.to_csv(output_file, index=False, header=False)
        return output_file
    except Exception as e:
        return f"Error in CSV to TXT conversion: {str(e)}"

# Handle file upload and conversion process
if uploaded_file:
    # Process uploaded file
    with NamedTemporaryFile(delete=False) as tmp_file:
        tmp_file.write(uploaded_file.getvalue())
        tmp_file_path = tmp_file.name

    result = "No conversion performed"  # Default value for result
    output_file = None  # Initialize output_file as None

    try:
        if conversion_type == "TXT to PDF":
            output_file = tmp_file_path.replace(".txt", ".pdf")
            result = convert_txt_to_pdf(tmp_file_path, output_file)
        elif conversion_type == "MD to PDF":
            output_file = tmp_file_path.replace(".md", ".pdf")
            result = convert_md_to_pdf(tmp_file_path, output_file)
        elif conversion_type == "HTML to PDF":
            output_file = tmp_file_path.replace(".html", ".pdf")
            result = convert_html_to_pdf(tmp_file_path, output_file)
        elif conversion_type == "EPUB to PDF":
            output_file = tmp_file_path.replace(".epub", ".pdf")
            result = convert_epub_to_pdf(tmp_file_path, output_file)
        elif conversion_type == "JSON to PDF":
            output_file = tmp_file_path.replace(".json", ".pdf")
            result = convert_json_to_pdf(tmp_file_path, output_file)
        elif conversion_type == "XML to PDF":
            output_file = tmp_file_path.replace(".xml", ".pdf")
            result = convert_xml_to_pdf(tmp_file_path, output_file)
        elif conversion_type == "QR Code Generator":
            output_file = tmp_file_path.replace(".txt", ".png")
            result = generate_qr_code(tmp_file_path, output_file)
        elif conversion_type == "Barcode Generator":
            output_file = tmp_file_path.replace(".txt", ".png")
            result = generate_barcode(tmp_file_path, output_file)
        elif conversion_type == "Image to PDF":
            output_file = tmp_file_path.replace((".jpg", ".jpeg", ".png"), ".pdf")
            result = convert_image_to_pdf(tmp_file_path, output_file)
        elif conversion_type == "Audio to Text":
            result = convert_audio_to_text(tmp_file_path)
        elif conversion_type == "Text to Speech":
            output_audio = tmp_file_path.replace(".txt", ".mp3")
            result = text_to_speech(tmp_file_path, output_audio)
        elif conversion_type == "Video to Audio":
            output_audio = tmp_file_path.replace(".mp4", ".mp3")
            result = convert_video_to_audio(tmp_file_path, output_audio)
        elif conversion_type == "Extract ZIP":
            output_folder = tmp_file_path.replace(".zip", "")
            result = extract_zip(tmp_file_path, output_folder)
        elif conversion_type == "Compress Folder":
            output_file = tmp_file_path + "_compressed"
            result = compress_folder(tmp_file_path, output_file)

        # Provide download link
        if output_file:
            st.success(f"Conversion completed successfully: {result}")
            st.download_button("Download Converted File", data=open(output_file, "rb").read(), file_name=output_file)
    except Exception as e:
        st.error(f"Error: {str(e)}")
